"""
Salem-Branded Reports API
Professional advisor/client-ready portfolio analysis reports
"""
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from models.schemas import (
    ReportData,
    ReportSummary,
    NarrativeBlock,
    MonteCarloBlock,
    StressScenarioResult,
    AssumptionsBlock,
    AppendixItem,
    KeyMetric,
    PercentilePathPoint,
    StressMetric,
    SuccessProbabilityPoint,
    TerminalWealthBucket,
    CashFlowProjection,
    IncomeSourcesTimeline,
    # Enhanced narrative report models
    NarrativeReportRequest,
    NarrativeReportResponse,
    EnhancedNarrativeReportModel,
    ExecutiveSummaryModel,
    IdentifiedRiskModel,
    RecommendationModel,
    FailureAnalysisModel,
    WorstCaseAnalysisModel,
    FailurePatternModel,
    WhatIfScenarioModel,
    RiskLevelEnum,
    RiskTypeEnum
)
from core.report_generator import (
    NarrativeEngine,
    RiskAnalyzer,
    RecommendationEngine,
    FailureAnalyzer,
    WorstCaseAnalyzer,
    RiskLevel,
    RiskType
)
from datetime import date, datetime
import logging
from typing import List
from io import BytesIO
import uuid
import os
import numpy as np

router = APIRouter()
logger = logging.getLogger(__name__)


def format_currency(value: float, decimals: int = 0) -> str:
    """Format value as currency"""
    if abs(value) >= 1_000_000:
        return f"${value/1_000_000:.{decimals}f}M"
    elif abs(value) >= 1_000:
        return f"${value/1_000:.{decimals}f}K"
    return f"${value:,.{decimals}f}"


def format_percent(value: float, decimals: int = 1) -> str:
    """Format value as percentage"""
    return f"{value * 100:.{decimals}f}%"


def generate_key_findings(
    success_prob: float,
    median_ending: float,
    starting: float,
    depletion_prob: float
) -> List[str]:
    """Generate contextual key findings"""
    findings = []
    
    # Success probability finding
    if success_prob >= 0.85:
        findings.append(
            f"Plan demonstrates strong {format_percent(success_prob)} probability of success, "
            f"indicating robust likelihood of meeting all financial objectives throughout the planning period."
        )
    elif success_prob >= 0.70:
        findings.append(
            f"Plan shows moderate {format_percent(success_prob)} probability of success. "
            f"While reasonable, consider stress testing and maintaining spending flexibility."
        )
    else:
        findings.append(
            f"Plan exhibits {format_percent(success_prob)} success probability, below recommended thresholds. "
            f"Adjustments to spending, allocation, or timeline recommended."
        )
    
    # Portfolio trajectory finding
    growth = (median_ending - starting) / starting
    if growth > 0.50:
        findings.append(
            f"Median scenario projects substantial portfolio growth to {format_currency(median_ending)}, "
            f"representing {format_percent(growth)} increase from starting value."
        )
    elif growth > 0:
        findings.append(
            f"Median scenario projects modest portfolio growth to {format_currency(median_ending)}."
        )
    else:
        findings.append(
            f"Median scenario projects controlled spend-down to {format_currency(median_ending)}, "
            f"aligned with decumulation objectives."
        )
    
    # Depletion risk finding
    if depletion_prob > 0.15:
        findings.append(
            f"Elevated depletion risk of {format_percent(depletion_prob)} warrants contingency planning "
            f"and potential adjustments to mitigate downside scenarios."
        )
    
    return findings


def generate_key_risks(
    success_prob: float,
    depletion_prob: float,
    equity_pct: float
) -> List[str]:
    """Generate key risk assessments"""
    risks = []
    
    if success_prob < 0.70:
        risks.append(
            "Current plan success probability falls below recommended thresholds, "
            "indicating material risk of not achieving stated financial objectives."
        )
    
    if depletion_prob > 0.20:
        risks.append(
            f"Portfolio depletion probability of {format_percent(depletion_prob)} represents "
            f"significant longevity risk requiring mitigation strategies."
        )
    
    if equity_pct > 0.80:
        risks.append(
            f"Aggressive {format_percent(equity_pct)} equity allocation increases exposure "
            f"to market volatility and sequence-of-returns risk, particularly critical during early retirement years."
        )
    elif equity_pct < 0.30:
        risks.append(
            f"Conservative {format_percent(equity_pct)} equity allocation may limit growth potential, "
            f"increasing exposure to inflation and longevity risk over extended time horizons."
        )
    
    risks.append(
        "Market conditions, inflation rates, and healthcare costs may differ materially from assumptions, "
        "affecting actual outcomes. Regular plan reviews recommended."
    )
    
    return risks


def generate_recommendations(
    success_prob: float,
    median_ending: float,
    starting: float
) -> List[str]:
    """Generate actionable recommendations"""
    recommendations = []
    
    if success_prob < 0.70:
        recommendations.extend([
            "Consider reducing discretionary spending by 10-15% to improve plan sustainability and increase success probability.",
            "Explore opportunities to defer retirement or incorporate part-time income to strengthen financial position.",
            "Review asset allocation with advisor to ensure appropriate risk/return profile for circumstances and objectives.",
            "Develop flexible spending contingencies to adapt to adverse market conditions if they occur."
        ])
    elif success_prob < 0.85:
        recommendations.extend([
            "Maintain current strategy with annual monitoring and periodic stress testing against various scenarios.",
            "Consider establishing spending flexibility thresholds to dynamically adjust to market performance.",
            "Review and update plan annually or following significant life events, market changes, or goal modifications.",
            "Evaluate tax optimization strategies including Roth conversions, strategic withdrawals, and loss harvesting."
        ])
    else:
        surplus = median_ending - starting
        if surplus > starting * 0.50:
            recommendations.append(
                "Plan demonstrates significant margin of success. Consider enhancing lifestyle spending, "
                "legacy goals, or charitable giving aligned with values and tax objectives."
            )
        recommendations.extend([
            "Explore tax-efficient wealth transfer strategies including lifetime gifting within annual exclusions.",
            "Review estate planning documents to ensure alignment with current intentions and tax law changes.",
            "Consider charitable giving opportunities through donor-advised funds or direct contributions.",
            "Maintain disciplined rebalancing and tax-loss harvesting within investment policy guidelines."
        ])
    
    return recommendations


@router.get("/reports/{plan_id}", response_model=ReportData)
async def get_report(plan_id: str):
    """
    Retrieve comprehensive Salem-branded portfolio analysis report.
    
    This endpoint generates a complete, advisor/client-ready report with:
    - Executive summary with key metrics
    - Narrative insights (findings, risks, recommendations)
    - Monte Carlo simulation results with percentile paths
    - Stress test scenario analysis
    - Planning assumptions and disclosures
    
    **Parameters:**
    - plan_id: Unique plan identifier (simulation run ID)
    
    **Returns:**
    - Complete ReportData structure ready for rendering
    """
    try:
        logger.info(f"Generating report for plan_id: {plan_id}")
        
        # TODO: In production, fetch from database/storage using plan_id
        # For now, generate sample report data demonstrating structure
        
        # Sample data representing a typical simulation result
        starting_portfolio = 4_500_000
        median_ending = 5_200_000
        p10_ending = 2_100_000
        p90_ending = 9_800_000
        success_prob = 0.82
        depletion_prob = 0.12
        equity_pct = 0.70
        years = 30
        num_runs = 1000
        
        # Generate key metrics for summary cards
        key_metrics = [
            KeyMetric(
                label="Success Probability",
                value=format_percent(success_prob),
                tooltip="Probability of meeting all spending needs throughout planning horizon",
                variant="success" if success_prob >= 0.85 else "warning" if success_prob >= 0.70 else "error"
            ),
            KeyMetric(
                label="Median Ending Portfolio",
                value=format_currency(median_ending),
                tooltip="50th percentile outcome - most likely portfolio value at end of planning period",
                variant="success"
            ),
            KeyMetric(
                label="10th Percentile (Downside)",
                value=format_currency(p10_ending),
                tooltip="Portfolio value in adverse scenarios - 90% of outcomes exceed this level",
                variant="neutral"
            ),
            KeyMetric(
                label="Depletion Risk",
                value=format_percent(depletion_prob),
                tooltip="Probability of portfolio depleting before end of planning horizon",
                variant="error" if depletion_prob > 0.15 else "warning"
            ),
        ]
        
        # Build report summary
        summary = ReportSummary(
            client_name="John & Mary Smith",  # TODO: Fetch from plan data
            scenario_name="Base Case Monte Carlo Analysis",
            as_of_date=date.today().isoformat(),
            advisor_name="Michael Chen, CFP®",  # TODO: Fetch from plan data
            firm_name="Salem Investment Counselors",
            key_metrics=key_metrics
        )
        
        # Generate narrative content
        narrative = NarrativeBlock(
            key_findings=generate_key_findings(success_prob, median_ending, starting_portfolio, depletion_prob),
            key_risks=generate_key_risks(success_prob, depletion_prob, equity_pct),
            recommendations=generate_recommendations(success_prob, median_ending, starting_portfolio)
        )
        
        # Build percentile path (yearly data points with 5 percentiles)
        percentile_path = []
        for year in range(years + 1):
            # Generate representative percentile values
            # In production, these would come from actual simulation results
            growth_factor_p50 = 1.06 ** year
            growth_factor_p5 = 0.92 ** year
            growth_factor_p10 = 0.98 ** year
            growth_factor_p25 = 1.02 ** year
            growth_factor_p75 = 1.10 ** year
            growth_factor_p90 = 1.12 ** year
            growth_factor_p95 = 1.14 ** year
            
            percentile_path.append(
                PercentilePathPoint(
                    year=year,
                    p5=starting_portfolio * growth_factor_p5,
                    p10=starting_portfolio * growth_factor_p10,
                    p25=starting_portfolio * growth_factor_p25,
                    p50=starting_portfolio * growth_factor_p50,
                    p75=starting_portfolio * growth_factor_p75,
                    p90=starting_portfolio * growth_factor_p90,
                    p95=starting_portfolio * growth_factor_p95
                )
            )
        
        # Generate success probability over time
        success_prob_over_time = []
        for year in range(1, years + 1):
            # Success probability tends to decrease over time as scenarios diverge
            year_success = success_prob * (0.98 ** (year / 10))  # Slight decay
            success_prob_over_time.append(
                SuccessProbabilityPoint(
                    year=year,
                    success_probability=min(1.0, year_success)
                )
            )
        
        # Generate terminal wealth distribution (histogram)
        terminal_wealth_dist = []
        bucket_ranges = [
            (0, 500_000, "$0-$500K"),
            (500_000, 1_000_000, "$500K-$1M"),
            (1_000_000, 2_000_000, "$1M-$2M"),
            (2_000_000, 3_000_000, "$2M-$3M"),
            (3_000_000, 5_000_000, "$3M-$5M"),
            (5_000_000, 7_500_000, "$5M-$7.5M"),
            (7_500_000, 10_000_000, "$7.5M-$10M"),
            (10_000_000, float('inf'), "$10M+"),
        ]
        
        # Simulate distribution (in production, use actual simulation results)
        import random
        random.seed(42)  # For reproducibility
        terminal_values = [
            median_ending * random.lognormvariate(0, 0.5) for _ in range(num_runs)
        ]
        
        for min_val, max_val, label in bucket_ranges:
            count = sum(1 for v in terminal_values if min_val <= v < max_val)
            if count > 0:
                terminal_wealth_dist.append(
                    TerminalWealthBucket(
                        bucket_label=label,
                        count=count,
                        min_value=min_val,
                        max_value=max_val if max_val != float('inf') else max(terminal_values),
                        percentage=count / num_runs
                    )
                )
        
        # Build Monte Carlo block
        monte_carlo = MonteCarloBlock(
            percentile_path=percentile_path,
            success_probability=success_prob,
            num_runs=num_runs,
            horizon_years=years,
            first_failure_year=None if success_prob > 0.50 else 22,
            success_probability_over_time=success_prob_over_time,
            terminal_wealth_distribution=terminal_wealth_dist
        )
        
        # Generate cash flow projections
        current_age = 48  # TODO: Get from plan data
        annual_spending = 240_000
        ss_income = 48_000
        ss_start_age = 67
        
        cash_flow_projections = []
        balance = starting_portfolio
        
        for year in range(1, years + 1):
            age = current_age + year
            
            # Income sources
            ss_this_year = ss_income if age >= ss_start_age else 0
            pension_this_year = 0  # TODO: Add pension logic
            total_income = ss_this_year + pension_this_year
            
            # Withdrawals (negative)
            withdrawals = -annual_spending * (1.03 ** year)  # Inflation adjusted
            
            # Taxes (simplified)
            taxable_income = max(0, -withdrawals + total_income)
            taxes = -taxable_income * 0.25
            
            # Investment return (simplified - use median growth)
            inv_return = balance * 0.06
            
            # Ending balance
            ending = balance + withdrawals + total_income + taxes + inv_return
            
            cash_flow_projections.append(
                CashFlowProjection(
                    year=year,
                    age=age,
                    beginning_balance=balance,
                    withdrawals=withdrawals,
                    income_sources_total=total_income,
                    taxes=taxes,
                    investment_return=inv_return,
                    ending_balance=ending
                )
            )
            
            balance = ending
        
        # Generate income timeline
        income_timeline = []
        for year in range(1, years + 1):
            age = current_age + year
            ss_this_year = ss_income if age >= ss_start_age else 0
            withdrawals_this_year = annual_spending * (1.03 ** year)
            
            income_timeline.append(
                IncomeSourcesTimeline(
                    year=year,
                    social_security=ss_this_year,
                    pension=0,  # TODO: Add pension
                    annuity=0,  # TODO: Add annuity
                    portfolio_withdrawals=withdrawals_this_year,
                    other_income=0
                )
            )
        
        # Generate stress test scenarios
        stress_tests = [
            StressScenarioResult(
                id="early_bear",
                name="Early Bear Market",
                description="Severe market downturn (-30% equity, -10% bonds) in first 3 years of plan",
                base_success_probability=success_prob,
                stressed_success_probability=max(0, success_prob - 0.15),
                base_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(success_prob), variant="success"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending), variant="neutral"),
                ],
                stressed_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(max(0, success_prob - 0.15)), variant="warning"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending * 0.75), variant="warning"),
                ],
                impact_severity="high"
            ),
            StressScenarioResult(
                id="high_inflation",
                name="Elevated Inflation",
                description="Persistent 5% annual inflation versus 3% base assumption throughout planning period",
                base_success_probability=success_prob,
                stressed_success_probability=max(0, success_prob - 0.10),
                base_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(success_prob), variant="success"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending), variant="neutral"),
                ],
                stressed_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(max(0, success_prob - 0.10)), variant="warning"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending * 0.82), variant="neutral"),
                ],
                impact_severity="medium"
            ),
            StressScenarioResult(
                id="lower_returns",
                name="Lower Market Returns",
                description="Market returns 2 percentage points below long-term historical averages across all asset classes",
                base_success_probability=success_prob,
                stressed_success_probability=max(0, success_prob - 0.18),
                base_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(success_prob), variant="success"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending), variant="neutral"),
                ],
                stressed_key_metrics=[
                    KeyMetric(label="Success Rate", value=format_percent(max(0, success_prob - 0.18)), variant="error"),
                    KeyMetric(label="Median Ending", value=format_currency(median_ending * 0.68), variant="warning"),
                ],
                impact_severity="high"
            ),
        ]
        
        # Build assumptions block
        assumptions = AssumptionsBlock(
            planning_horizon_years=years,
            real_return_mean=0.065,  # 6.5% nominal minus 3% inflation
            real_return_std=0.12,
            inflation_rate=0.03,
            spending_rule_description="Fixed dollar spending adjusted annually for inflation (3% rate)",
            other_assumptions={
                "equity_allocation": format_percent(equity_pct),
                "fixed_income_allocation": format_percent(0.25),
                "cash_allocation": format_percent(0.05),
                "starting_portfolio": format_currency(starting_portfolio),
                "annual_spending": format_currency(240_000),  # Example: $20K/month
                "social_security": "$48,000/year starting age 67",
                "healthcare_costs": "$15,000/year starting age 65, inflating at 5%",
                "tax_rate": "25% marginal rate",
                "estate_goal": format_currency(2_000_000)
            }
        )
        
        # Build appendix sections
        appendix = [
            AppendixItem(
                title="Methodology",
                content=[
                    "Monte Carlo simulation utilizing geometric Brownian motion to model portfolio returns and volatility.",
                    "Asset class returns modeled with log-normal distribution incorporating historical mean returns and standard deviations.",
                    f"{num_runs:,} independent simulation paths generated to capture range of potential outcomes.",
                    "Spending rules, income sources, taxes, and healthcare costs incorporated into each scenario.",
                    "Success defined as maintaining positive portfolio balance throughout entire planning horizon."
                ]
            ),
            AppendixItem(
                title="Important Limitations",
                content=[
                    "Simulations based on historical return and volatility assumptions which may not reflect future market conditions.",
                    "Extreme market events ('black swans') may not be adequately captured in historical data.",
                    "Actual tax treatment depends on individual circumstances and may differ from assumptions.",
                    "Inflation assumptions may not reflect actual cost increases, particularly for healthcare.",
                    "Does not constitute investment advice or guarantee of future results. Consult with qualified professionals."
                ]
            ),
            AppendixItem(
                title="Regular Review Recommended",
                content=[
                    "Financial plans should be reviewed annually or following significant life events.",
                    "Market conditions, tax laws, and personal circumstances change over time.",
                    "Adjust assumptions and strategies as needed to maintain plan relevance and accuracy.",
                    "Contact your Salem Investment Counselors advisor to discuss updates or questions."
                ]
            ),
        ]
        
        # Assemble complete report
        report = ReportData(
            report_id=plan_id,
            summary=summary,
            narrative=narrative,
            monte_carlo=monte_carlo,
            stress_tests=stress_tests,
            assumptions=assumptions,
            appendix=appendix,
            cash_flow_projection=cash_flow_projections,
            income_timeline=income_timeline
        )
        
        logger.info(f"Report generated successfully for plan_id: {plan_id}")
        return report
        
    except Exception as e:
        logger.error(f"Failed to generate report for plan_id {plan_id}: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Report generation failed: {str(e)}")


@router.post("/reports/{plan_id}/export/powerpoint")
async def export_powerpoint(plan_id: str):
    """
    Export portfolio analysis report as PowerPoint presentation (.pptx).
    
    Generates a professional, Salem-branded PowerPoint presentation with:
    - Executive summary slide with key metrics
    - Monte Carlo analysis results
    - Stress test scenarios
    - Planning assumptions
    - Recommendations
    
    **Parameters:**
    - plan_id: Unique plan identifier (simulation run ID)
    
    **Returns:**
    - PowerPoint file (.pptx) for download
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        logger.info(f"Generating PowerPoint export for plan_id: {plan_id}")
        
        # Get report data (reuse existing report generation logic)
        report = await get_report(plan_id)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define Salem brand colors
        SALEM_NAVY = RGBColor(0, 51, 93)  # #00335D
        SALEM_GOLD = RGBColor(180, 151, 89)  # #B49759
        DARK_GRAY = RGBColor(51, 51, 51)
        LIGHT_GRAY = RGBColor(245, 245, 245)
        
        # ==========================================
        # SLIDE 1: Title Slide
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SALEM_NAVY
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Portfolio Analysis Report"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Client name
        client_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
        client_frame = client_box.text_frame
        client_frame.text = report.summary.client_name
        client_para = client_frame.paragraphs[0]
        client_para.font.size = Pt(32)
        client_para.font.color.rgb = SALEM_GOLD
        client_para.alignment = PP_ALIGN.CENTER
        
        # Date and firm
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{report.summary.firm_name} | {report.summary.as_of_date}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = RGBColor(200, 200, 200)
        footer_para.alignment = PP_ALIGN.CENTER
        
        # ==========================================
        # SLIDE 2: Executive Summary
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        # Key Metrics Grid (2x2)
        metrics_positions = [
            (0.5, 1.2, 4.25, 1.5),  # Top left
            (5.25, 1.2, 4.25, 1.5),  # Top right
            (0.5, 2.9, 4.25, 1.5),  # Bottom left
            (5.25, 2.9, 4.25, 1.5),  # Bottom right
        ]
        
        for i, metric in enumerate(report.summary.key_metrics[:4]):
            if i >= len(metrics_positions):
                break
                
            left, top, width, height = metrics_positions[i]
            
            # Metric box with border
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_GRAY
            shape.line.color.rgb = SALEM_NAVY
            shape.line.width = Pt(2)
            
            # Metric label
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 0.2), 
                Inches(width - 0.4), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = metric.label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = DARK_GRAY
            label_para.font.bold = True
            
            # Metric value
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 0.7), 
                Inches(width - 0.4), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = metric.value
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            
            # Color based on variant
            if metric.variant == "success":
                value_para.font.color.rgb = RGBColor(16, 185, 129)  # Green
            elif metric.variant == "warning":
                value_para.font.color.rgb = RGBColor(245, 158, 11)  # Amber
            elif metric.variant == "error":
                value_para.font.color.rgb = RGBColor(239, 68, 68)  # Red
            else:
                value_para.font.color.rgb = SALEM_NAVY
        
        # ==========================================
        # SLIDE 3: Key Findings
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Key Findings"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        # Findings list
        findings_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
        findings_frame = findings_box.text_frame
        findings_frame.word_wrap = True
        
        for i, finding in enumerate(report.narrative.key_findings):
            if i > 0:
                findings_frame.add_paragraph()
            p = findings_frame.paragraphs[i]
            p.text = f"• {finding}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(16)
            p.level = 0
        
        # ==========================================
        # SLIDE 4: Key Risks
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Key Risks & Considerations"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        # Risks list
        risks_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
        risks_frame = risks_box.text_frame
        risks_frame.word_wrap = True
        
        for i, risk in enumerate(report.narrative.key_risks):
            if i > 0:
                risks_frame.add_paragraph()
            p = risks_frame.paragraphs[i]
            p.text = f"• {risk}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(16)
            p.level = 0
        
        # ==========================================
        # SLIDE 5: Recommendations
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Recommendations"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        # Recommendations list
        rec_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True
        
        for i, rec in enumerate(report.narrative.recommendations):
            if i > 0:
                rec_frame.add_paragraph()
            p = rec_frame.paragraphs[i]
            p.text = f"• {rec}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(16)
            p.level = 0
        
        # ==========================================
        # SLIDE 6: Stress Tests Summary
        # ==========================================
        if report.stress_tests:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = "Stress Test Analysis"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = SALEM_NAVY
            
            # Stress tests
            y_pos = 1.3
            for stress in report.stress_tests[:3]:  # Show top 3
                # Scenario name
                name_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.4)
                )
                name_frame = name_box.text_frame
                name_frame.text = stress.name
                name_para = name_frame.paragraphs[0]
                name_para.font.size = Pt(18)
                name_para.font.bold = True
                name_para.font.color.rgb = SALEM_NAVY
                
                # Description
                desc_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y_pos + 0.4), Inches(8.4), Inches(0.5)
                )
                desc_frame = desc_box.text_frame
                desc_frame.text = stress.description
                desc_frame.word_wrap = True
                desc_para = desc_frame.paragraphs[0]
                desc_para.font.size = Pt(13)
                desc_para.font.color.rgb = DARK_GRAY
                
                # Impact
                impact_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y_pos + 0.95), Inches(8.4), Inches(0.3)
                )
                impact_frame = impact_box.text_frame
                success_change = stress.stressed_success_probability - stress.base_success_probability
                impact_frame.text = f"Success Probability Impact: {success_change:+.1%}"
                impact_para = impact_frame.paragraphs[0]
                impact_para.font.size = Pt(14)
                impact_para.font.bold = True
                
                if stress.impact_severity == "high":
                    impact_para.font.color.rgb = RGBColor(239, 68, 68)  # Red
                elif stress.impact_severity == "medium":
                    impact_para.font.color.rgb = RGBColor(245, 158, 11)  # Amber
                else:
                    impact_para.font.color.rgb = RGBColor(16, 185, 129)  # Green
                
                y_pos += 1.8
        
        # ==========================================
        # SLIDE 7: Planning Assumptions
        # ==========================================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Planning Assumptions"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        # Assumptions
        assumptions_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
        assumptions_frame = assumptions_box.text_frame
        assumptions_frame.word_wrap = True
        
        assumptions_list = [
            f"Planning Horizon: {report.assumptions.planning_horizon_years} years",
            f"Inflation Rate: {report.assumptions.inflation_rate:.1%}",
            f"Expected Return: {report.assumptions.real_return_mean:.1%} (real)",
            f"Return Volatility: {report.assumptions.real_return_std:.1%}",
            f"Spending Rule: {report.assumptions.spending_rule_description}",
        ]
        
        # Add other assumptions
        for key, value in list(report.assumptions.other_assumptions.items())[:5]:
            formatted_key = key.replace('_', ' ').title()
            assumptions_list.append(f"{formatted_key}: {value}")
        
        for i, assumption in enumerate(assumptions_list):
            if i > 0:
                assumptions_frame.add_paragraph()
            p = assumptions_frame.paragraphs[i]
            p.text = f"• {assumption}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
        
        # ==========================================
        # Save to bytes buffer
        # ==========================================
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Generate filename
        client_name = report.summary.client_name.replace(' ', '_')
        filename = f"Portfolio_Analysis_{client_name}_{date.today().strftime('%Y%m%d')}.pptx"
        
        logger.info(f"PowerPoint generated successfully for plan_id: {plan_id}")
        
        # Return as streaming response
        return StreamingResponse(
            pptx_buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to generate PowerPoint for plan_id {plan_id}: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PowerPoint generation failed: {str(e)}")


@router.post("/reports/export/pdf")
async def export_pdf_from_results(simulation_results: dict):
    """
    Export portfolio analysis report as professional PDF document using actual simulation results.
    
    Accepts simulation results from frontend and generates a comprehensive, Salem-branded PDF report.
    
    **Parameters:**
    - simulation_results: Complete simulation data including metrics, paths, and inputs
    
    **Returns:**
    - PDF file for download
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, Image as RLImage, KeepTogether
        )
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import HexColor
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend
        import matplotlib.pyplot as plt
        import numpy as np
        
        logger.info(f"Generating PDF export from simulation results")
        
        # Extract data from simulation results
        metrics = simulation_results.get('metrics', {})
        percentile_paths = simulation_results.get('percentile_paths', [])
        inputs = simulation_results.get('inputs', {})
        client_info = simulation_results.get('client_info', {})
        
        # Create PDF buffer
        pdf_buffer = BytesIO()
        
        # Define Salem brand colors
        SALEM_NAVY = HexColor('#00335D')
        SALEM_GOLD = HexColor('#B49759')
        SALEM_GREEN = HexColor('#4B8F29')
        SALEM_RED = HexColor('#9E2A2B')
        DARK_GRAY = HexColor('#333333')
        LIGHT_GRAY = HexColor('#F5F5F5')
        
        # Create document
        client_name = client_info.get('client_name', 'Client')
        doc = SimpleDocTemplate(
            pdf_buffer,
            pagesize=letter,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch,
            title=f"Portfolio Analysis - {client_name}",
            author="Salem Investment Counselors"
        )
        
        # Custom styles
        styles = getSampleStyleSheet()
        
        # Helper function to add or get style
        def add_style_if_not_exists(name, **kwargs):
            if name not in styles:
                styles.add(ParagraphStyle(name=name, **kwargs))
            return styles[name]
        
        # Cover title style
        add_style_if_not_exists(
            'CoverTitle',
            parent=styles['Heading1'],
            fontSize=32,
            textColor=SALEM_NAVY,
            spaceAfter=12,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Cover subtitle style
        add_style_if_not_exists(
            'CoverSubtitle',
            parent=styles['Normal'],
            fontSize=18,
            textColor=SALEM_GOLD,
            spaceAfter=6,
            alignment=TA_CENTER,
            fontName='Helvetica'
        )
        
        # Section heading style
        add_style_if_not_exists(
            'SectionHeading',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=SALEM_NAVY,
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold',
            borderWidth=0,
            borderColor=SALEM_NAVY,
            borderPadding=0,
            leftIndent=0
        )
        
        # Subsection heading
        add_style_if_not_exists(
            'SubHeading',
            parent=styles['Heading2'],
            fontSize=13,
            textColor=SALEM_NAVY,
            spaceAfter=8,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        )
        
        # Body text
        add_style_if_not_exists(
            'BodyText',
            parent=styles['Normal'],
            fontSize=10,
            textColor=DARK_GRAY,
            spaceAfter=8,
            alignment=TA_JUSTIFY,
            fontName='Helvetica'
        )
        
        # Bullet style
        add_style_if_not_exists(
            'Bullet',
            parent=styles['Normal'],
            fontSize=10,
            textColor=DARK_GRAY,
            spaceAfter=6,
            leftIndent=20,
            fontName='Helvetica'
        )
        
        # Build PDF content
        story = []
        
        # ==========================================
        # COVER PAGE
        # ==========================================
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph("Portfolio Analysis Report", styles['CoverTitle']))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(client_name, styles['CoverSubtitle']))
        story.append(Paragraph("Monte Carlo Retirement Analysis", styles['CoverSubtitle']))
        story.append(Spacer(1, 0.5*inch))
        
        # Firm info
        firm_style = ParagraphStyle(
            'FirmInfo',
            parent=styles['Normal'],
            fontSize=12,
            textColor=DARK_GRAY,
            alignment=TA_CENTER
        )
        story.append(Paragraph("<b>Salem Investment Counselors</b>", firm_style))
        story.append(Paragraph(f"As of: {date.today().strftime('%B %d, %Y')}", firm_style))
        
        story.append(PageBreak())
        
        # ==========================================
        # EXECUTIVE SUMMARY
        # ==========================================
        story.append(Paragraph("Executive Summary", styles['SectionHeading']))
        story.append(Spacer(1, 0.2*inch))
        
        # Extract key metrics
        success_prob = metrics.get('success_probability', 0)
        median_ending = metrics.get('ending_median', 0)
        p10_ending = metrics.get('ending_p10', 0)
        p90_ending = metrics.get('ending_p90', 0)
        depletion_prob = metrics.get('depletion_probability', 0)
        
        # Key metrics table
        metrics_data = [['Metric', 'Value', 'Assessment']]
        
        # Success Probability
        success_assessment = '✓ Strong' if success_prob >= 0.85 else '⚠ Adequate' if success_prob >= 0.70 else '✗ At Risk'
        metrics_data.append(['Success Probability', format_percent(success_prob, 1), success_assessment])
        
        # Median Ending
        metrics_data.append(['Median Ending Portfolio', format_currency(median_ending, 0), 'Neutral'])
        
        # 10th Percentile
        metrics_data.append(['10th Percentile (Downside)', format_currency(p10_ending, 0), 'Neutral'])
        
        # Depletion Risk
        depl_assessment = '✓ Low' if depletion_prob < 0.10 else '⚠ Moderate' if depletion_prob < 0.20 else '✗ High'
        metrics_data.append(['Depletion Risk', format_percent(depletion_prob, 1), depl_assessment])
        
        metrics_table = Table(metrics_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
        metrics_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), SALEM_NAVY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('TOPPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), LIGHT_GRAY),
            ('GRID', (0, 0), (-1, -1), 0.5, DARK_GRAY),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, LIGHT_GRAY]),
        ]))
        story.append(metrics_table)
        story.append(Spacer(1, 0.3*inch))
        
        # ==========================================
        # MONTE CARLO ANALYSIS - Create wealth fan chart with ACTUAL DATA
        # ==========================================
        story.append(Paragraph("Monte Carlo Simulation Results", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        # Create wealth fan chart from actual percentile paths
        if percentile_paths and len(percentile_paths) > 0:
            fig, ax = plt.subplots(figsize=(7, 4.5))
            
            years = [p.get('year', i) for i, p in enumerate(percentile_paths)]
            p5 = [p.get('p5', 0) / 1_000_000 for p in percentile_paths]
            p10 = [p.get('p10', 0) / 1_000_000 for p in percentile_paths]
            p25 = [p.get('p25', 0) / 1_000_000 for p in percentile_paths]
            p50 = [p.get('median', p.get('p50', 0)) / 1_000_000 for p in percentile_paths]
            p75 = [p.get('p75', 0) / 1_000_000 for p in percentile_paths]
            p90 = [p.get('p90', 0) / 1_000_000 for p in percentile_paths]
            p95 = [p.get('p95', 0) / 1_000_000 for p in percentile_paths]
            
            # Plot percentile bands
            ax.fill_between(years, p5, p95, alpha=0.15, color='#00335D', label='5th-95th %ile')
            ax.fill_between(years, p10, p90, alpha=0.2, color='#00335D', label='10th-90th %ile')
            ax.fill_between(years, p25, p75, alpha=0.3, color='#00335D', label='25th-75th %ile')
            ax.plot(years, p50, color='#00335D', linewidth=2.5, label='Median (50th %ile)')
            
            ax.set_xlabel('Year', fontsize=11, fontweight='bold')
            ax.set_ylabel('Portfolio Value ($M)', fontsize=11, fontweight='bold')
            ax.set_title('Projected Wealth Outcomes', fontsize=13, fontweight='bold', color='#00335D')
            ax.legend(loc='best', fontsize=9, framealpha=0.9)
            ax.grid(True, alpha=0.3, linestyle='--')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            # Save chart to buffer
            chart_buffer = BytesIO()
            plt.tight_layout()
            plt.savefig(chart_buffer, format='png', dpi=150, bbox_inches='tight')
            plt.close()
            chart_buffer.seek(0)
            
            # Add chart to PDF
            chart_image = RLImage(chart_buffer, width=6*inch, height=3.75*inch)
            story.append(chart_image)
            story.append(Spacer(1, 0.2*inch))
        
        # Monte Carlo summary text
        num_runs = inputs.get('num_simulations', 1000)
        years_horizon = inputs.get('years_in_retirement', 30)
        
        mc_summary = f"""
        This Monte Carlo simulation ran {num_runs:,} scenarios over a 
        {years_horizon}-year planning horizon. The median projection shows 
        the most likely outcome, while the shaded bands represent the range of potential results. 
        The plan demonstrates a {format_percent(success_prob, 0)} probability 
        of successfully meeting all financial objectives throughout the planning period.
        """
        story.append(Paragraph(mc_summary.strip(), styles['BodyText']))
        
        story.append(PageBreak())
        
        # ==========================================
        # PLANNING ASSUMPTIONS
        # ==========================================
        story.append(Paragraph("Planning Assumptions", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        portfolio_value = inputs.get('portfolio_value', 0)
        monthly_spending = abs(inputs.get('monthly_spending', 0))
        annual_spending = monthly_spending * 12
        equity_return = inputs.get('equity_return_annual', 0)
        fi_return = inputs.get('fi_return_annual', 0)
        inflation = inputs.get('inflation_annual', 0)
        
        assumptions_text = f"""
        <b>Starting Portfolio:</b> {format_currency(portfolio_value, 0)}<br/>
        <b>Annual Spending:</b> {format_currency(annual_spending, 0)} ({format_currency(monthly_spending, 0)}/month)<br/>
        <b>Planning Horizon:</b> {years_horizon} years<br/>
        <b>Equity Return:</b> {format_percent(equity_return, 1)}<br/>
        <b>Fixed Income Return:</b> {format_percent(fi_return, 1)}<br/>
        <b>Inflation Rate:</b> {format_percent(inflation, 1)}<br/>
        <b>Simulation Runs:</b> {num_runs:,}<br/>
        """
        story.append(Paragraph(assumptions_text, styles['BodyText']))
        story.append(Spacer(1, 0.3*inch))
        
        # ==========================================
        # KEY INSIGHTS
        # ==========================================
        story.append(Paragraph("Key Insights", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        # Generate insights based on actual data
        insights = []
        
        # Success probability insight
        if success_prob >= 0.85:
            insights.append(
                f"Plan demonstrates strong {format_percent(success_prob, 0)} probability of success, "
                f"indicating robust likelihood of meeting all financial objectives."
            )
        elif success_prob >= 0.70:
            insights.append(
                f"Plan shows moderate {format_percent(success_prob, 0)} probability of success. "
                f"Consider stress testing and maintaining spending flexibility."
            )
        else:
            insights.append(
                f"Plan exhibits {format_percent(success_prob, 0)} success probability, below recommended thresholds. "
                f"Adjustments to spending, allocation, or timeline recommended."
            )
        
        # Portfolio trajectory
        growth = (median_ending - portfolio_value) / portfolio_value if portfolio_value > 0 else 0
        if growth > 0.50:
            insights.append(
                f"Median scenario projects substantial portfolio growth to {format_currency(median_ending, 0)}, "
                f"representing {format_percent(growth, 0)} increase from starting value."
            )
        elif growth > 0:
            insights.append(
                f"Median scenario projects modest portfolio growth to {format_currency(median_ending, 0)}."
            )
        else:
            insights.append(
                f"Median scenario projects controlled spend-down to {format_currency(median_ending, 0)}, "
                f"aligned with decumulation objectives."
            )
        
        # Depletion risk
        if depletion_prob > 0.15:
            insights.append(
                f"Elevated depletion risk of {format_percent(depletion_prob, 0)} warrants contingency planning "
                f"and potential adjustments to mitigate downside scenarios."
            )
        
        for insight in insights:
            story.append(Paragraph(f"• {insight}", styles['Bullet']))
        
        story.append(PageBreak())
        
        # ==========================================
        # APPENDIX
        # ==========================================
        story.append(Paragraph("Important Disclosures", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        disclosures = [
            "This analysis is based on Monte Carlo simulation utilizing historical return and volatility assumptions.",
            "Past performance does not guarantee future results. Actual outcomes may differ materially from projections.",
            "Simulations model thousands of potential scenarios but cannot predict actual market conditions or returns.",
            "Tax treatment depends on individual circumstances and may differ from assumptions used in this analysis.",
            "This report is for informational purposes only and does not constitute investment advice.",
            "Consult with your advisor before making any financial decisions based on this analysis."
        ]
        
        for disclosure in disclosures:
            story.append(Paragraph(f"• {disclosure}", styles['Bullet']))
        
        # ==========================================
        # Build PDF
        # ==========================================
        def add_page_number(canvas, doc):
            """Add page numbers to each page"""
            page_num = canvas.getPageNumber()
            text = f"Page {page_num}"
            canvas.saveState()
            canvas.setFont('Helvetica', 9)
            canvas.setFillColor(DARK_GRAY)
            canvas.drawRightString(7.5*inch, 0.5*inch, text)
            canvas.drawString(0.75*inch, 0.5*inch, f"Salem Investment Counselors | {client_name}")
            canvas.restoreState()
        
        doc.build(story, onFirstPage=add_page_number, onLaterPages=add_page_number)
        pdf_buffer.seek(0)
        
        # Generate filename
        safe_client_name = client_name.replace(' ', '_')
        filename = f"Portfolio_Analysis_{safe_client_name}_{date.today().strftime('%Y%m%d')}.pdf"
        
        logger.info(f"PDF generated successfully with actual simulation data")
        
        # Return as streaming response
        return StreamingResponse(
            pdf_buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to generate PDF from simulation results: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {str(e)}")


@router.post("/reports/export/powerpoint")
async def export_powerpoint_from_results(simulation_results: dict):
    """
    Export PowerPoint presentation using actual simulation results.
    Accepts simulation data directly from frontend.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        
        # Extract actual data from passed simulation results
        metrics = simulation_results.get('metrics', {})
        percentile_paths = simulation_results.get('percentile_paths', [])
        inputs = simulation_results.get('inputs', {})
        client_info = simulation_results.get('client_info', {})
        
        # Extract key metrics
        success_prob = metrics.get('success_probability', 0)
        median_ending = metrics.get('ending_median', 0)
        p10_ending = metrics.get('ending_p10', 0)
        p90_ending = metrics.get('ending_p90', 0)
        depletion_prob = metrics.get('depletion_probability', 0)
        
        # Extract input parameters
        starting_value = inputs.get('portfolio_value', 0)
        monthly_spending = inputs.get('monthly_spending', 0)
        annual_spending = monthly_spending * 12
        years_retirement = inputs.get('years_in_retirement', 30)
        expected_return = inputs.get('expected_return', 0.07)
        volatility = inputs.get('volatility', 0.12)
        inflation = inputs.get('inflation_rate', 0.025)
        
        # Client information
        client_name = client_info.get('client_name', 'Client')
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define Salem brand colors
        SALEM_NAVY = RGBColor(0, 51, 93)
        SALEM_GOLD = RGBColor(180, 151, 89)
        DARK_GRAY = RGBColor(51, 51, 51)
        LIGHT_GRAY = RGBColor(245, 245, 245)
        
        # ========== SLIDE 1: Title Slide ==========
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SALEM_NAVY
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Portfolio Analysis Report"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        client_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
        client_frame = client_box.text_frame
        client_frame.text = client_name
        client_para = client_frame.paragraphs[0]
        client_para.font.size = Pt(32)
        client_para.font.color.rgb = SALEM_GOLD
        client_para.alignment = PP_ALIGN.CENTER
        
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Salem Investment Counselors | {date.today().strftime('%B %d, %Y')}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = RGBColor(200, 200, 200)
        footer_para.alignment = PP_ALIGN.CENTER
        
        # ========== SLIDE 2: Executive Summary ==========
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        metrics_positions = [
            (0.5, 1.2, 4.25, 1.5),
            (5.25, 1.2, 4.25, 1.5),
            (0.5, 2.9, 4.25, 1.5),
            (5.25, 2.9, 4.25, 1.5),
        ]
        
        actual_metrics = [
            ("Success Probability", f"{success_prob:.1%}"),
            ("Median Ending Value", f"${median_ending / 1_000_000:.2f}M"),
            ("10th Percentile", f"${p10_ending / 1_000_000:.2f}M"),
            ("Depletion Risk", f"{depletion_prob:.1%}"),
        ]
        
        for i, (label, value) in enumerate(actual_metrics):
            left, top, width, height = metrics_positions[i]
            
            shape = slide.shapes.add_shape(
                1, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_GRAY
            shape.line.color.rgb = SALEM_NAVY
            shape.line.width = Pt(2)
            
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 0.2), 
                Inches(width - 0.4), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = DARK_GRAY
            label_para.font.bold = True
            
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 0.7), 
                Inches(width - 0.4), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = value
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = SALEM_NAVY
            value_para.alignment = PP_ALIGN.CENTER
        
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.2))
        note_frame = note_box.text_frame
        note_frame.word_wrap = True
        
        if success_prob >= 0.85:
            assessment = "strong probability of success"
        elif success_prob >= 0.70:
            assessment = "adequate probability of success"
        else:
            assessment = "elevated risk that may require plan adjustments"
            
        note_text = f"Based on {years_retirement}-year Monte Carlo simulation with 10,000 trials, "
        note_text += f"this portfolio demonstrates a {assessment}. "
        note_text += f"Starting with ${starting_value / 1_000_000:.2f}M and spending ${annual_spending:,.0f} annually, "
        note_text += f"the median projection reaches ${median_ending / 1_000_000:.2f}M."
        
        note_frame.text = note_text
        for para in note_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = DARK_GRAY
            para.space_before = Pt(6)
            para.space_after = Pt(6)
        
        # ========== SLIDE 3: Wealth Chart ==========
        if percentile_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = "Portfolio Wealth Projection"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = SALEM_NAVY
            
            fig, ax = plt.subplots(figsize=(10, 5))
            
            years = [p.get('year', i) for i, p in enumerate(percentile_paths)]
            p5 = [p.get('p5', 0) / 1_000_000 for p in percentile_paths]
            p10 = [p.get('p10', 0) / 1_000_000 for p in percentile_paths]
            p25 = [p.get('p25', 0) / 1_000_000 for p in percentile_paths]
            p50 = [p.get('median', p.get('p50', 0)) / 1_000_000 for p in percentile_paths]
            p75 = [p.get('p75', 0) / 1_000_000 for p in percentile_paths]
            p90 = [p.get('p90', 0) / 1_000_000 for p in percentile_paths]
            p95 = [p.get('p95', 0) / 1_000_000 for p in percentile_paths]
            
            ax.fill_between(years, p5, p95, alpha=0.1, color='#00335D', label='5th-95th Percentile')
            ax.fill_between(years, p10, p90, alpha=0.2, color='#00335D', label='10th-90th Percentile')
            ax.fill_between(years, p25, p75, alpha=0.3, color='#00335D', label='25th-75th Percentile')
            ax.plot(years, p50, color='#B49759', linewidth=3, label='Median')
            
            ax.set_xlabel('Year', fontsize=12, fontweight='bold')
            ax.set_ylabel('Portfolio Value ($M)', fontsize=12, fontweight='bold')
            ax.set_title('Wealth Projection', fontsize=14, fontweight='bold', color='#00335D')
            ax.legend(loc='upper left')
            ax.grid(True, alpha=0.3)
            
            chart_path = f"/tmp/pptx_chart_{uuid.uuid4().hex}.png"
            plt.tight_layout()
            plt.savefig(chart_path, dpi=150, bbox_inches='tight')
            plt.close()
            
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), width=Inches(9))
            
            if os.path.exists(chart_path):
                os.remove(chart_path)
        
        # ========== SLIDE 4: Assumptions ==========
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Analysis Assumptions"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = SALEM_NAVY
        
        params_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        params_frame = params_box.text_frame
        params_frame.word_wrap = True
        
        params_text = f"""Starting Portfolio: ${starting_value:,.0f}

Annual Spending: ${annual_spending:,.0f}

Time Horizon: {years_retirement} years

Expected Return: {expected_return:.1%}

Volatility: {volatility:.1%}

Inflation: {inflation:.1%}

Method: Monte Carlo (10,000 trials)"""
        
        params_frame.text = params_text
        for para in params_frame.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = DARK_GRAY
            para.space_after = Pt(12)
        
        # Save to buffer
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return StreamingResponse(
            pptx_buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename=Portfolio_Analysis_{client_name.replace(' ', '_')}_{date.today().strftime('%Y%m%d')}.pptx"
            }
        )
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint from results: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PowerPoint generation failed: {str(e)}")


@router.post("/reports/{plan_id}/export/pdf")
async def export_pdf(plan_id: str):
    """
    Export portfolio analysis report as professional PDF document.
    
    Generates a comprehensive, Salem-branded PDF report with:
    - Cover page with branding
    - Executive summary with key metrics
    - Monte Carlo wealth fan chart
    - Success probability trends
    - Stress test analysis
    - Cash flow projections
    - Planning assumptions
    - Professional formatting and styling
    
    **Parameters:**
    - plan_id: Unique plan identifier (simulation run ID)
    
    **Returns:**
    - PDF file for download
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, Image as RLImage, KeepTogether
        )
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import HexColor
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend
        import matplotlib.pyplot as plt
        import numpy as np
        
        logger.info(f"Generating PDF export for plan_id: {plan_id}")
        
        # Get report data
        report = await get_report(plan_id)
        
        # Create PDF buffer
        pdf_buffer = BytesIO()
        
        # Define Salem brand colors
        SALEM_NAVY = HexColor('#00335D')
        SALEM_GOLD = HexColor('#B49759')
        SALEM_GREEN = HexColor('#4B8F29')
        SALEM_RED = HexColor('#9E2A2B')
        DARK_GRAY = HexColor('#333333')
        LIGHT_GRAY = HexColor('#F5F5F5')
        
        # Create document
        doc = SimpleDocTemplate(
            pdf_buffer,
            pagesize=letter,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch,
            title=f"Portfolio Analysis - {report.summary.client_name}",
            author="Salem Investment Counselors"
        )
        
        # Custom styles
        styles = getSampleStyleSheet()
        
        # Helper function to add or get style
        def add_style_if_not_exists(name, **kwargs):
            if name not in styles:
                styles.add(ParagraphStyle(name=name, **kwargs))
            return styles[name]
        
        # Cover title style
        add_style_if_not_exists(
            'CoverTitle',
            parent=styles['Heading1'],
            fontSize=32,
            textColor=SALEM_NAVY,
            spaceAfter=12,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Cover subtitle style
        add_style_if_not_exists(
            'CoverSubtitle',
            parent=styles['Normal'],
            fontSize=18,
            textColor=SALEM_GOLD,
            spaceAfter=6,
            alignment=TA_CENTER,
            fontName='Helvetica'
        )
        
        # Section heading style
        add_style_if_not_exists(
            'SectionHeading',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=SALEM_NAVY,
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold',
            borderWidth=0,
            borderColor=SALEM_NAVY,
            borderPadding=0,
            leftIndent=0
        )
        
        # Subsection heading
        add_style_if_not_exists(
            'SubHeading',
            parent=styles['Heading2'],
            fontSize=13,
            textColor=SALEM_NAVY,
            spaceAfter=8,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        )
        
        # Body text
        add_style_if_not_exists(
            'BodyText',
            parent=styles['Normal'],
            fontSize=10,
            textColor=DARK_GRAY,
            spaceAfter=8,
            alignment=TA_JUSTIFY,
            fontName='Helvetica'
        )
        
        # Bullet style
        add_style_if_not_exists(
            'Bullet',
            parent=styles['Normal'],
            fontSize=10,
            textColor=DARK_GRAY,
            spaceAfter=6,
            leftIndent=20,
            fontName='Helvetica'
        )
        
        # Build PDF content
        story = []
        
        # ==========================================
        # COVER PAGE
        # ==========================================
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph("Portfolio Analysis Report", styles['CoverTitle']))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(report.summary.client_name, styles['CoverSubtitle']))
        story.append(Paragraph(report.summary.scenario_name, styles['CoverSubtitle']))
        story.append(Spacer(1, 0.5*inch))
        
        # Firm info
        firm_style = ParagraphStyle(
            'FirmInfo',
            parent=styles['Normal'],
            fontSize=12,
            textColor=DARK_GRAY,
            alignment=TA_CENTER
        )
        story.append(Paragraph(f"<b>{report.summary.firm_name}</b>", firm_style))
        story.append(Paragraph(f"Prepared by: {report.summary.advisor_name}", firm_style))
        story.append(Paragraph(f"As of: {report.summary.as_of_date}", firm_style))
        
        story.append(PageBreak())
        
        # ==========================================
        # EXECUTIVE SUMMARY
        # ==========================================
        story.append(Paragraph("Executive Summary", styles['SectionHeading']))
        story.append(Spacer(1, 0.2*inch))
        
        # Key metrics table
        metrics_data = [['Metric', 'Value', 'Assessment']]
        for metric in report.summary.key_metrics:
            assessment = '✓ Strong' if metric.variant == 'success' else \
                        '⚠ Caution' if metric.variant == 'warning' else \
                        '✗ At Risk' if metric.variant == 'error' else 'Neutral'
            metrics_data.append([metric.label, metric.value, assessment])
        
        metrics_table = Table(metrics_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
        metrics_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), SALEM_NAVY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('TOPPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), LIGHT_GRAY),
            ('GRID', (0, 0), (-1, -1), 0.5, DARK_GRAY),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, LIGHT_GRAY]),
        ]))
        story.append(metrics_table)
        story.append(Spacer(1, 0.3*inch))
        
        # ==========================================
        # MONTE CARLO ANALYSIS - Create wealth fan chart
        # ==========================================
        story.append(Paragraph("Monte Carlo Simulation Results", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        # Create wealth fan chart
        fig, ax = plt.subplots(figsize=(7, 4.5))
        
        years = [p.year for p in report.monte_carlo.percentile_path]
        p5 = [p.p5 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p10 = [p.p10 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p25 = [p.p25 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p50 = [p.p50 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p75 = [p.p75 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p90 = [p.p90 / 1_000_000 for p in report.monte_carlo.percentile_path]
        p95 = [p.p95 / 1_000_000 for p in report.monte_carlo.percentile_path]
        
        # Plot percentile bands
        ax.fill_between(years, p5, p95, alpha=0.15, color='#00335D', label='5th-95th %ile')
        ax.fill_between(years, p10, p90, alpha=0.2, color='#00335D', label='10th-90th %ile')
        ax.fill_between(years, p25, p75, alpha=0.3, color='#00335D', label='25th-75th %ile')
        ax.plot(years, p50, color='#00335D', linewidth=2.5, label='Median (50th %ile)')
        
        ax.set_xlabel('Year', fontsize=11, fontweight='bold')
        ax.set_ylabel('Portfolio Value ($M)', fontsize=11, fontweight='bold')
        ax.set_title('Projected Wealth Outcomes', fontsize=13, fontweight='bold', color='#00335D')
        ax.legend(loc='best', fontsize=9, framealpha=0.9)
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        # Save chart to buffer
        chart_buffer = BytesIO()
        plt.tight_layout()
        plt.savefig(chart_buffer, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        chart_buffer.seek(0)
        
        # Add chart to PDF
        chart_image = RLImage(chart_buffer, width=6*inch, height=3.75*inch)
        story.append(chart_image)
        story.append(Spacer(1, 0.2*inch))
        
        # Monte Carlo summary text
        mc_summary = f"""
        This Monte Carlo simulation ran {report.monte_carlo.num_runs:,} scenarios over a 
        {report.monte_carlo.horizon_years}-year planning horizon. The median projection shows 
        the most likely outcome, while the shaded bands represent the range of potential results. 
        The plan demonstrates a {format_percent(report.monte_carlo.success_probability, 0)} probability 
        of successfully meeting all financial objectives throughout the planning period.
        """
        story.append(Paragraph(mc_summary.strip(), styles['BodyText']))
        
        story.append(PageBreak())
        
        # ==========================================
        # KEY FINDINGS
        # ==========================================
        story.append(Paragraph("Key Findings", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        for finding in report.narrative.key_findings:
            story.append(Paragraph(f"• {finding}", styles['Bullet']))
        
        story.append(Spacer(1, 0.3*inch))
        
        # ==========================================
        # KEY RISKS
        # ==========================================
        story.append(Paragraph("Key Risks & Considerations", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        for risk in report.narrative.key_risks:
            story.append(Paragraph(f"• {risk}", styles['Bullet']))
        
        story.append(Spacer(1, 0.3*inch))
        
        # ==========================================
        # RECOMMENDATIONS
        # ==========================================
        story.append(Paragraph("Recommendations", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        for i, rec in enumerate(report.narrative.recommendations, 1):
            story.append(Paragraph(f"{i}. {rec}", styles['Bullet']))
        
        story.append(PageBreak())
        
        # ==========================================
        # STRESS TEST ANALYSIS
        # ==========================================
        if report.stress_tests:
            story.append(Paragraph("Stress Test Analysis", styles['SectionHeading']))
            story.append(Spacer(1, 0.1*inch))
            
            stress_intro = """
            The following stress tests evaluate plan resilience under adverse scenarios. 
            Each test compares baseline assumptions against stressed conditions to quantify potential impacts.
            """
            story.append(Paragraph(stress_intro.strip(), styles['BodyText']))
            story.append(Spacer(1, 0.15*inch))
            
            for stress in report.stress_tests[:3]:
                # Stress test header
                story.append(Paragraph(f"<b>{stress.name}</b>", styles['SubHeading']))
                story.append(Paragraph(stress.description, styles['BodyText']))
                story.append(Spacer(1, 0.1*inch))
                
                # Comparison table
                stress_data = [
                    ['Scenario', 'Success Probability', 'Impact'],
                    ['Base Case', format_percent(stress.base_success_probability, 1), '—'],
                    ['Stressed', format_percent(stress.stressed_success_probability, 1), 
                     format_percent(stress.stressed_success_probability - stress.base_success_probability, 1)]
                ]
                
                stress_table = Table(stress_data, colWidths=[2*inch, 1.75*inch, 1.5*inch])
                stress_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), LIGHT_GRAY),
                    ('TEXTCOLOR', (0, 0), (-1, 0), DARK_GRAY),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 8),
                    ('GRID', (0, 0), (-1, -1), 0.5, DARK_GRAY),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                ]))
                story.append(stress_table)
                story.append(Spacer(1, 0.2*inch))
            
            story.append(PageBreak())
        
        # ==========================================
        # CASH FLOW PROJECTION (sample years)
        # ==========================================
        if report.cash_flow_projection:
            story.append(Paragraph("Cash Flow Projection", styles['SectionHeading']))
            story.append(Spacer(1, 0.1*inch))
            
            # Show first 10 years
            cf_data = [['Year', 'Age', 'Beginning', 'Withdrawals', 'Income', 'Taxes', 'Return', 'Ending']]
            for cf in report.cash_flow_projection[:10]:
                cf_data.append([
                    str(cf.year),
                    str(cf.age),
                    format_currency(cf.beginning_balance, 0),
                    format_currency(cf.withdrawals, 0),
                    format_currency(cf.income_sources_total, 0),
                    format_currency(cf.taxes, 0),
                    format_currency(cf.investment_return, 0),
                    format_currency(cf.ending_balance, 0)
                ])
            
            cf_table = Table(cf_data, colWidths=[0.4*inch, 0.4*inch, 0.85*inch, 0.85*inch, 
                                                  0.75*inch, 0.65*inch, 0.75*inch, 0.85*inch])
            cf_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), SALEM_NAVY),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'RIGHT'),
                ('ALIGN', (0, 0), (1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 8),
                ('FONTSIZE', (0, 1), (-1, -1), 7),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                ('TOPPADDING', (0, 0), (-1, 0), 8),
                ('TOPPADDING', (0, 1), (-1, -1), 5),
                ('BOTTOMPADDING', (0, 1), (-1, -1), 5),
                ('GRID', (0, 0), (-1, -1), 0.5, DARK_GRAY),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, LIGHT_GRAY]),
            ]))
            story.append(cf_table)
            story.append(Spacer(1, 0.1*inch))
            story.append(Paragraph(
                "<i>Note: Table shows first 10 years. Full projection available upon request.</i>",
                styles['BodyText']
            ))
            
            story.append(PageBreak())
        
        # ==========================================
        # PLANNING ASSUMPTIONS
        # ==========================================
        story.append(Paragraph("Planning Assumptions", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        assumptions_text = f"""
        <b>Planning Horizon:</b> {report.assumptions.planning_horizon_years} years<br/>
        <b>Expected Return (Real):</b> {format_percent(report.assumptions.real_return_mean, 1)}<br/>
        <b>Return Volatility:</b> {format_percent(report.assumptions.real_return_std, 1)}<br/>
        <b>Inflation Rate:</b> {format_percent(report.assumptions.inflation_rate, 1)}<br/>
        <b>Spending Rule:</b> {report.assumptions.spending_rule_description}<br/>
        """
        story.append(Paragraph(assumptions_text, styles['BodyText']))
        story.append(Spacer(1, 0.15*inch))
        
        # Additional assumptions
        if report.assumptions.other_assumptions:
            story.append(Paragraph("<b>Additional Assumptions:</b>", styles['SubHeading']))
            for key, value in list(report.assumptions.other_assumptions.items())[:8]:
                formatted_key = key.replace('_', ' ').title()
                story.append(Paragraph(f"• <b>{formatted_key}:</b> {value}", styles['Bullet']))
        
        story.append(PageBreak())
        
        # ==========================================
        # APPENDIX
        # ==========================================
        story.append(Paragraph("Appendix", styles['SectionHeading']))
        story.append(Spacer(1, 0.1*inch))
        
        for item in report.appendix:
            story.append(Paragraph(f"<b>{item.title}</b>", styles['SubHeading']))
            for content_item in item.content:
                story.append(Paragraph(f"• {content_item}", styles['Bullet']))
            story.append(Spacer(1, 0.15*inch))
        
        # ==========================================
        # Build PDF
        # ==========================================
        def add_page_number(canvas, doc):
            """Add page numbers to each page"""
            page_num = canvas.getPageNumber()
            text = f"Page {page_num}"
            canvas.saveState()
            canvas.setFont('Helvetica', 9)
            canvas.setFillColor(DARK_GRAY)
            canvas.drawRightString(7.5*inch, 0.5*inch, text)
            canvas.drawString(0.75*inch, 0.5*inch, f"Salem Investment Counselors | {report.summary.client_name}")
            canvas.restoreState()
        
        doc.build(story, onFirstPage=add_page_number, onLaterPages=add_page_number)
        pdf_buffer.seek(0)
        
        # Generate filename
        client_name = report.summary.client_name.replace(' ', '_')
        filename = f"Portfolio_Analysis_{client_name}_{date.today().strftime('%Y%m%d')}.pdf"
        
        logger.info(f"PDF generated successfully for plan_id: {plan_id}")
        
        # Return as streaming response
        return StreamingResponse(
            pdf_buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to generate PDF for plan_id {plan_id}: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {str(e)}")


# ============================================================================
# ENHANCED NARRATIVE REPORT ENDPOINT
# ============================================================================

@router.post("/narrative", response_model=NarrativeReportResponse)
async def generate_narrative_report(request: NarrativeReportRequest):
    """
    Generate enhanced narrative report with risk analysis and recommendations.
    
    Produces client-ready reports with:
    - Executive summary in plain English
    - Top identified risks with severity and mitigation strategies
    - Prioritized actionable recommendations
    - Optional failure analysis (what scenarios fail, when, why)
    - Optional worst-case analysis (10th percentile deep-dive)
    
    Args:
        request: NarrativeReportRequest with simulation results and parameters
    
    Returns:
        NarrativeReportResponse with complete narrative report
    """
    try:
        logger.info("Generating enhanced narrative report")
        
        # Initialize engines
        narrative_engine = NarrativeEngine()
        risk_analyzer = RiskAnalyzer()
        recommendation_engine = RecommendationEngine()
        
        # 1. Generate Executive Summary
        executive_summary = narrative_engine.generate_executive_summary(
            success_probability=request.success_probability,
            median_ending_value=request.median_ending_value,
            percentile_10_value=request.percentile_10_value,
            percentile_90_value=request.percentile_90_value,
            starting_portfolio=request.starting_portfolio,
            years_to_model=request.years_to_model,
            current_age=request.current_age,
            monthly_spending=request.monthly_spending,
            has_goals=request.has_goals,
            goals_on_track_count=request.goals_on_track_count,
            total_goals=request.total_goals
        )
        
        exec_summary_model = ExecutiveSummaryModel(
            plan_overview=executive_summary.plan_overview,
            success_probability_narrative=executive_summary.success_probability_narrative,
            key_strengths=executive_summary.key_strengths,
            key_concerns=executive_summary.key_concerns,
            bottom_line=executive_summary.bottom_line
        )
        
        # 2. Identify Risks
        risks = risk_analyzer.identify_risks(
            success_probability=request.success_probability,
            median_ending=request.median_ending_value,
            percentile_10=request.percentile_10_value,
            failure_scenarios=np.array([]),  # Simplified - use all_paths if provided
            starting_portfolio=request.starting_portfolio,
            annual_spending=abs(request.monthly_spending) * 12,
            years_to_model=request.years_to_model,
            current_age=request.current_age,
            horizon_age=request.current_age + request.years_to_model,
            equity_pct=request.equity_pct,
            monthly_spending=request.monthly_spending
        )
        
        # Convert to API models
        risk_models = [
            IdentifiedRiskModel(
                risk_type=_convert_risk_type_to_enum(risk.risk_type),
                severity=_convert_risk_level_to_enum(risk.severity),
                probability=risk.probability,
                potential_impact=risk.potential_impact,
                description=risk.description,
                mitigation_strategy=risk.mitigation_strategy,
                priority_rank=risk.priority_rank
            )
            for risk in risks
        ]
        
        # 3. Generate Recommendations
        recommendations = recommendation_engine.generate_recommendations(
            risks=risks,
            success_probability=request.success_probability,
            starting_portfolio=request.starting_portfolio,
            annual_spending=abs(request.monthly_spending) * 12,
            equity_pct=request.equity_pct,
            years_to_model=request.years_to_model
        )
        
        recommendation_models = [
            RecommendationModel(
                title=rec.title,
                description=rec.description,
                expected_benefit=rec.expected_benefit,
                implementation_steps=rec.implementation_steps,
                priority=rec.priority,
                category=rec.category
            )
            for rec in recommendations
        ]
        
        # 4. Failure Analysis (if requested and paths provided)
        failure_analysis_model = None
        if request.include_failure_analysis and request.all_paths:
            try:
                failure_analyzer = FailureAnalyzer()
                paths_array = np.array(request.all_paths)
                
                failure_analysis = failure_analyzer.analyze_failures(
                    all_paths=paths_array,
                    success_threshold=0,  # Portfolio depletion threshold
                    years_to_model=request.years_to_model,
                    starting_portfolio=request.starting_portfolio,
                    annual_spending=abs(request.monthly_spending) * 12
                )
                
                # Convert patterns to models
                pattern_models = [
                    FailurePatternModel(
                        pattern_name=pattern.pattern_name,
                        frequency=pattern.frequency,
                        typical_failure_year=pattern.typical_failure_year,
                        description=pattern.description,
                        prevention_strategy=pattern.prevention_strategy
                    )
                    for pattern in failure_analysis['patterns']
                ]
                
                failure_analysis_model = FailureAnalysisModel(
                    failure_count=failure_analysis['failure_count'],
                    failure_rate=failure_analysis['failure_rate'],
                    avg_failure_year=failure_analysis.get('avg_failure_year'),
                    median_failure_year=failure_analysis.get('median_failure_year'),
                    earliest_failure_year=failure_analysis.get('earliest_failure_year'),
                    patterns=pattern_models,
                    summary=failure_analysis['summary'],
                    prevention_strategies=failure_analysis['prevention_strategies']
                )
            except Exception as e:
                logger.warning(f"Failed to generate failure analysis: {str(e)}")
        
        # 5. Worst-Case Analysis (if requested and paths provided)
        worst_case_model = None
        if request.include_worst_case_analysis and request.all_paths:
            try:
                worst_case_analyzer = WorstCaseAnalyzer()
                paths_array = np.array(request.all_paths)
                
                worst_case = worst_case_analyzer.analyze_worst_case(
                    all_paths=paths_array,
                    percentile_10_value=request.percentile_10_value,
                    starting_portfolio=request.starting_portfolio,
                    annual_spending=abs(request.monthly_spending) * 12,
                    years_to_model=request.years_to_model
                )
                
                # Convert what-if scenarios
                what_if_models = [
                    WhatIfScenarioModel(
                        scenario=scenario['scenario'],
                        change=scenario['change'],
                        impact=scenario['impact'],
                        trade_off=scenario['trade_off']
                    )
                    for scenario in worst_case['what_if_scenarios']
                ]
                
                worst_case_model = WorstCaseAnalysisModel(
                    percentile_10_value=worst_case['percentile_10_value'],
                    max_drawdown_pct=worst_case['max_drawdown_pct'],
                    turning_point_year=worst_case['turning_point_year'],
                    recovery_time_years=worst_case['recovery_time_years'],
                    description=worst_case['description'],
                    recovery_strategies=worst_case['recovery_strategies'],
                    what_if_scenarios=what_if_models
                )
            except Exception as e:
                logger.warning(f"Failed to generate worst-case analysis: {str(e)}")
        
        # Assemble complete report
        report = EnhancedNarrativeReportModel(
            executive_summary=exec_summary_model,
            identified_risks=risk_models,
            recommendations=recommendation_models,
            failure_analysis=failure_analysis_model,
            worst_case_analysis=worst_case_model,
            report_generated_at=datetime.utcnow().isoformat()
        )
        
        logger.info(
            f"Report generated: {len(risk_models)} risks, "
            f"{len(recommendation_models)} recommendations"
        )
        
        return NarrativeReportResponse(
            report=report,
            success=True,
            message="Narrative report generated successfully"
        )
        
    except ValueError as e:
        logger.error(f"Validation error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Error generating narrative report: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate narrative report: {str(e)}"
        )


def _convert_risk_type_to_enum(risk_type: RiskType) -> RiskTypeEnum:
    """Convert engine RiskType to API enum"""
    mapping = {
        RiskType.SEQUENCE_OF_RETURNS: RiskTypeEnum.SEQUENCE_OF_RETURNS,
        RiskType.LONGEVITY: RiskTypeEnum.LONGEVITY,
        RiskType.INFLATION: RiskTypeEnum.INFLATION,
        RiskType.HEALTHCARE_COSTS: RiskTypeEnum.HEALTHCARE_COSTS,
        RiskType.PORTFOLIO_DEPLETION: RiskTypeEnum.PORTFOLIO_DEPLETION,
        RiskType.TAX_INEFFICIENCY: RiskTypeEnum.TAX_INEFFICIENCY,
        RiskType.SPENDING_UNSUSTAINABLE: RiskTypeEnum.SPENDING_UNSUSTAINABLE,
        RiskType.CONCENTRATION: RiskTypeEnum.CONCENTRATION,
        RiskType.MARKET_VOLATILITY: RiskTypeEnum.MARKET_VOLATILITY,
    }
    return mapping[risk_type]


def _convert_risk_level_to_enum(risk_level: RiskLevel) -> RiskLevelEnum:
    """Convert engine RiskLevel to API enum"""
    mapping = {
        RiskLevel.LOW: RiskLevelEnum.LOW,
        RiskLevel.MODERATE: RiskLevelEnum.MODERATE,
        RiskLevel.HIGH: RiskLevelEnum.HIGH,
        RiskLevel.CRITICAL: RiskLevelEnum.CRITICAL,
    }
    return mapping[risk_level]
